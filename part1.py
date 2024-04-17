import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

def create_presentation_with_text_and_images(text, images):
    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Add a slide to the presentation for text
    slide_text = presentation.slides.add_slide(presentation.slide_layouts[5])  # Choose a slide layout
    text_box_text = slide_text.shapes.add_textbox(left=0, top=0, width=presentation.slide_width, height=presentation.slide_height)
    text_frame_text = text_box_text.text_frame
    text_frame_text.text = text

    # Add slides to the presentation for each image
    for image in images:
        slide_image = presentation.slides.add_slide(presentation.slide_layouts[5])  # Choose a slide layout
        left = top = Inches(1)
    
# Resize image (assuming you want to fit within remaining slide area)
    image_width = presentation.slide_width - left * 2  # Account for left & right margins
    image_height = presentation.slide_height - top * 2  # Account for top & bottom margins
    slide_image.shapes.add_picture(image, left, top, width=image_width, height=image_height)

    return presentation

def extract_pdf_text_and_images(pdf_file_path):
    doc = fitz.open(pdf_file_path)
    full_text = ""
    images = []

    for page in doc:
        full_text += page.get_text()

        # Extract images
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_path = f"image{img_index}.png"  # Save the image to a file
            with open(image_path, "wb") as image_file:
                image_file.write(image_bytes)
            images.append(image_path)

    doc.close()
    return full_text, images

# Example usage
pdf_file_path = "c:\\Users\\LENOVO\\Desktop\\New folder\\aaa.pdf"
extracted_text, extracted_images = extract_pdf_text_and_images(pdf_file_path)

# Create a PowerPoint presentation with the extracted text and images
presentation = create_presentation_with_text_and_images(extracted_text, extracted_images)

# Save the presentation to a file
presentation.save("extracted_text_and_images_presentation.pptx")
